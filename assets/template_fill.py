#!/usr/bin/env python3
"""
template_fill.py — 把内容填进「你自己的 .pptx 模板」，外观完全沿用模板（借鉴 ppt-master）。
运行：/usr/bin/python3 template_fill.py <模板.pptx> <内容.json> -o <输出.pptx>

模板里用占位符标记要填的地方（标完一次，以后复用）：
  · 文本占位符：在任意文本框 / 表格单元格里写 {{key}}，会被 content["key"] 替换（保留原格式）。
  · 重复表格行：在模板表格的「模板行」第一个单元格写 {{#each LISTKEY}}，该行其余单元格用列名占位
      （如 {{name}} {{amount}}）。会按 content["LISTKEY"]（dict 列表）克隆此行逐条填充。
  · 图片占位：给模板里的图片设「替代文字(alt text)」为 {{img:KEY}}，会用 content["img"][KEY]（本地图片路径）
      替换该图片，保持原位置与尺寸。

内容 JSON 示例：
  {
    "title": "复旦科创 2026 中期回顾",
    "subtitle": "Pipeline 投资分析",
    "projects": [
      {"name":"清纯半导体","sector":"半导体","amount":"3,000","status":"预立项"},
      {"name":"鹿明机器人","sector":"具身智能","amount":"2,000","status":"协议签署"}
    ],
    "img": {"cover": "/path/to/cover.jpg"}
  }
设计原则：只换内容、不动版式；占位符替换保留原 run 的字体/字号/颜色。
"""
import sys, json, re, argparse, copy
from pptx import Presentation
from pptx.util import Emu
from pptx.oxml.ns import qn

TOKEN = re.compile(r"\{\{\s*([^}]+?)\s*\}\}")
EACH = re.compile(r"\{\{\s*#each\s+([A-Za-z0-9_]+)\s*\}\}")


def _replace_in_paragraph(para, mapping):
    """在一个段落内替换 {{key}}，尽量保留 run 格式。"""
    # 先逐 run 替换（token 完整落在单个 run 的常见情况，完美保留格式）
    changed = False
    for run in para.runs:
        if "{{" in run.text:
            new = TOKEN.sub(lambda m: str(mapping.get(m.group(1).strip(), m.group(0))), run.text)
            if new != run.text:
                run.text = new; changed = True
    # 若 token 跨 run 拆分：合并整段文本到首 run，清空其余（保留首 run 格式）
    full = "".join(r.text for r in para.runs)
    if "{{" in full and TOKEN.search(full):
        new = TOKEN.sub(lambda m: str(mapping.get(m.group(1).strip(), m.group(0))), full)
        if new != full and para.runs:
            para.runs[0].text = new
            for r in para.runs[1:]:
                r.text = ""
            changed = True
    return changed


def _fill_text_frame(tf, mapping):
    for para in tf.paragraphs:
        _replace_in_paragraph(para, mapping)


def _fill_table(tbl, content):
    # 1) 找重复行模板：某行任一单元格含 {{#each KEY}}
    template_row_idx, list_key = None, None
    for ri, row in enumerate(tbl.rows):
        for cell in row.cells:
            m = EACH.search(cell.text)
            if m:
                template_row_idx, list_key = ri, m.group(1); break
        if list_key:
            break

    if list_key and list_key in content:
        items = content[list_key] or []
        tr_template = tbl.rows[template_row_idx]._tr
        tbl_el = tbl._tbl
        # 逐条克隆模板行（插到模板行后），按列名 {{col}} 填充
        anchor = tr_template
        for item in items:
            new_tr = copy.deepcopy(tr_template)
            tbl_el.insert(list(tbl_el).index(anchor) + 1, new_tr)
            anchor = new_tr
        # 删除原模板行
        tbl_el.remove(tr_template)
        # 现在重新遍历克隆出来的行填充（克隆行替换了模板行位置起）
        # 重新构建：用 python-pptx 重新读取行
        # 把 {{#each KEY}} 标记从每行清掉，并按 item 填充
        # 注意：克隆顺序与 items 对应
        clone_rows = [r for r in tbl.rows][template_row_idx: template_row_idx + len(items)]
        for row, item in zip(clone_rows, items):
            for cell in row.cells:
                # 去掉 #each 标记
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.text = EACH.sub("", run.text)
                _fill_text_frame(cell.text_frame, item)
    else:
        # 普通表格：仅做标量 token 替换
        for row in tbl.rows:
            for cell in row.cells:
                _fill_text_frame(cell.text_frame, content)


def _fill_images(slide, content):
    imgmap = content.get("img") or {}
    if not imgmap:
        return
    for shape in list(slide.shapes):
        try:
            is_pic = (shape.shape_type == 13)  # MSO_SHAPE_TYPE.PICTURE
        except Exception:
            is_pic = False
        if is_pic:
            descr = ""
            cn = shape._element.find(".//" + qn("p:cNvPr"))
            if cn is not None:
                descr = cn.get("descr") or ""
            m = re.search(r"\{\{\s*img:([A-Za-z0-9_]+)\s*\}\}", descr or "")
            if m and m.group(1) in imgmap:
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape._element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(imgmap[m.group(1)], left, top, width, height)


def fill(template_path, content, out_path):
    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                _fill_text_frame(shape.text_frame, content)
            if shape.has_table:
                _fill_table(shape.table, content)
        _fill_images(slide, content)
    prs.save(out_path)
    return out_path


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("template"); ap.add_argument("content")
    ap.add_argument("-o", "--out", required=True)
    a = ap.parse_args()
    with open(a.content, encoding="utf-8") as f:
        content = json.load(f)
    out = fill(a.template, content, a.out)
    print("WROTE", out)


if __name__ == "__main__":
    main()
