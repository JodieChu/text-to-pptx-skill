import glob
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

imgs = sorted(glob.glob("/Users/jodiechu/ai_projects/.tools-cache/img-cache/*.jpg") + glob.glob("/Users/jodiechu/ai_projects/.tools-cache/img-cache/*.png"))
placeholder_img = imgs[0]

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BRAND = RGBColor(0x0B, 0x3D, 0x2E)  # 品牌墨绿（模拟某机构标准色）
GOLD = RGBColor(0xC8, 0x90, 0x1C)
blank = prs.slide_layouts[6]

# ---- 封面 ----
s = prs.slides.add_slide(blank)
# 品牌色条（左竖条）—— 设计元素，填充时应原样保留
bar = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.35), Inches(7.5))
bar.fill.solid(); bar.fill.fore_color.rgb = BRAND; bar.line.fill.background()
# 标题
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.6), Inches(8.5), Inches(1.2)); tf = tb.text_frame
p = tf.paragraphs[0]; r = p.add_run(); r.text = "{{title}}"; r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = BRAND
# 副标题
tb2 = s.shapes.add_textbox(Inches(0.9), Inches(3.9), Inches(8.5), Inches(0.8)); tf2 = tb2.text_frame
p2 = tf2.paragraphs[0]; r2 = p2.add_run(); r2.text = "{{subtitle}}"; r2.font.size = Pt(20); r2.font.color.rgb = GOLD
# 封面图占位（右侧），alt text 标记
pic = s.shapes.add_picture(placeholder_img, Inches(9.6), Inches(0.0), Inches(3.733), Inches(7.5))
cn = pic._element.find(".//" + qn("p:cNvPr")); cn.set("descr", "{{img:cover}}")

# ---- 表格页 ----
s2 = prs.slides.add_slide(blank)
tb3 = s2.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11), Inches(0.8)); tf3 = tb3.text_frame
pr = tf3.paragraphs[0]; rr = pr.add_run(); rr.text = "{{table_title}}"; rr.font.size = Pt(28); rr.font.bold = True; rr.font.color.rgb = BRAND
rows, cols = 2, 4
gtbl = s2.shapes.add_table(rows, cols, Inches(0.6), Inches(1.6), Inches(12.1), Inches(0.9)).table
heads = ["项目名称", "所属赛道", "拟投(万)", "当前状态"]
for c, h in enumerate(heads):
    cell = gtbl.cell(0, c); cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = BRAND
    rp = cell.text_frame.paragraphs[0]; rp.runs[0].font.color.rgb = RGBColor(0xFF,0xFF,0xFF); rp.runs[0].font.bold = True; rp.runs[0].font.size = Pt(13)
# 模板行（带 #each + 列名占位）
tmpl = ["{{#each projects}}{{name}}", "{{sector}}", "{{amount}}", "{{status}}"]
for c, t in enumerate(tmpl):
    cell = gtbl.cell(1, c); cell.text = t
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)

prs.save("sample_template.pptx")
print("WROTE sample_template.pptx | 占位图:", placeholder_img.split("/")[-1])
